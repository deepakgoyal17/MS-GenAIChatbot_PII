from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a new PowerPoint presentation
prs = Presentation()

# Define title and content slide layout
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Slide 1: Title
slide1 = prs.slides.add_slide(title_slide_layout)
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "Production Support Overview"
subtitle.text = "Shift Staffing Matrix, Roles & Responsibilities"

# Slide 2: Roles & Responsibilities
slide2 = prs.slides.add_slide(content_slide_layout)
slide2.shapes.title.text = "Production Support Roles & Responsibilities"
slide2.placeholders[1].text = (
    "🔹 Team Lead\n"
    "- Manage incidents, coordinate tasks, review SOPs, lead shift handovers\n\n"
    "🔹 Production Manager\n"
    "- Own SLAs, manage risks, compliance, resource planning, stakeholder comms\n\n"
    "🔹 Product Analyst\n"
    "- Validate reports, resolve data issues, support users, liaise with business"
)

# Slide 3: Key Responsibilities
slide3 = prs.slides.add_slide(content_slide_layout)
slide3.shapes.title.text = "Key Responsibilities Overview"
slide3.placeholders[1].text = (
    "🔸 Incident Management\n"
    "🔸 Infra & ETL Monitoring\n"
    "🔸 User Query Handling\n"
    "🔸 Change & Release Mgmt\n"
    "🔸 Risk & Compliance Checks\n"
    "🔸 Report Validation & QA\n"
    "🔸 Documentation & SOPs\n"
    "🔸 Performance Tuning\n"
    "🔸 Shift Coordination & Handover"
)

# Slide 4: Shift Staffing Matrix
slide4 = prs.slides.add_slide(content_slide_layout)
slide4.shapes.title.text = "Shift Staffing Matrix (India vs UK)"
slide4.placeholders[1].text = (
    "✅ India Shift\n"
    "- ETL Monitoring, Infra, Risk Checks, Early QA\n"
    "- Suggested: 4–5 members\n\n"
    "✅ UK Shift\n"
    "- User Support, Report Validation, Deployments, Biz Comms\n"
    "- Suggested: 3–4 members\n\n"
    "🟢 0.5 FTE overlap for smooth handover"
)

# Slide 5: Decision Matrix Summary
slide5 = prs.slides.add_slide(content_slide_layout)
slide5.shapes.title.text = "Decision Matrix (Simplified View)"
slide5.placeholders[1].text = (
    "🔴 High Workload: ETL Monitoring, Incident Mgmt, User Queries\n"
    "🟠 Medium: Infra Monitoring, Report QA\n"
    "🟡 Low–Medium: Change Mgmt, Business Comms\n"
    "🟢 Low: Risk Checks, Shift Handover\n\n"
    "📌 Use volume & SLA data to fine-tune staffing"
)

# Save presentation
prs.save("Production_Support_Shift_Staffing_Presentation.pptx")
